from pptx import Presentation
import os


def create_presentation(report_text):

    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = (
        "Autonomous AI Research Report"
    )

    slide.placeholders[1].text = (
        "Generated Automatically using AI"
    )

    # Create slides
    sections = report_text.split("\n\n")

    for section in sections[:8]:

        if len(section.strip()) > 20:

            layout = prs.slide_layouts[1]

            slide = prs.slides.add_slide(layout)

            lines = section.split("\n")

            title = lines[0][:50]

            content = "\n".join(lines[1:])[:1000]

            slide.shapes.title.text = title

            slide.placeholders[1].text = content

    os.makedirs("presentations", exist_ok=True)

    prs.save(
        "presentations/final_research_presentation.pptx"
    )


if __name__ == "__main__":

    report_path = "reports/final_research_report.txt"

    with open(report_path, "r", encoding="utf-8") as f:
        report_text = f.read()

    print("Generating PowerPoint presentation...")

    create_presentation(report_text)

    print("Presentation created successfully!")